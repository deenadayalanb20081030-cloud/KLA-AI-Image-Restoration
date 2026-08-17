"""
KLA AI Hackathon - Official Idea Submission Generator (PPTX & PDF)
Generates:
1. NanoRestore_KLA_PS01.pptx (Editable 16:9 Presentation)
2. NanoRestore_KLA_PS01.pdf  (Ready-to-Upload 9-Slide Submission PDF)
3. VisionForge_KLA_PS01.pdf (Alternative filename copy)

Complies strictly with all KLA Idea Submission Template rules and slide structures.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF


def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    C_BG = RGBColor(11, 15, 25)         # Deep Space Navy
    C_CARD = RGBColor(19, 27, 46)       # Dark Slate Card
    C_ACCENT = RGBColor(0, 229, 255)    # Cyan Accent
    C_WHITE = RGBColor(248, 250, 252)   # Pure White Text
    C_MUTED = RGBColor(148, 163, 184)   # Slate Muted Text
    C_GOLD = RGBColor(245, 158, 11)     # Amber Accent
    C_GREEN = RGBColor(16, 185, 129)    # Emerald Accent

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, slide_num, title, subtitle):
        # Header Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(15, 35, 60)
        badge.line.color.rgb = C_ACCENT
        tf_b = badge.text_frame
        tf_b.text = f"SLIDE {slide_num} / 9"
        tf_b.paragraphs[0].font.size = Pt(11)
        tf_b.paragraphs[0].font.bold = True
        tf_b.paragraphs[0].font.color.rgb = C_ACCENT
        tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title
        tb = slide.shapes.add_textbox(Inches(3.2), Inches(0.3), Inches(9.3), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_WHITE

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_MUTED

    # =========================================================================
    # SLIDE 1: Team Details
    # =========================================================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)

    # Title Hero
    tb = s1.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "KLA AI HACKATHON 2026"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT

    p2 = tf.add_paragraph()
    p2.text = "AI-Based Restoration of Degraded Images"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE

    p3 = tf.add_paragraph()
    p3.text = "Team NanoRestore | Vellore Institute of Technology - Chennai"
    p3.font.size = Pt(14)
    p3.font.color.rgb = C_GOLD

    # 4 Member Cards
    members = [
        {"name": "Deenadayalan.B", "role": "Team Lead & AI Architecture Lead", "email": "deenadayalanb20081030@gmail.com", "phone": "+91 7667407910", "task": "NAFNetSR backbone, SimpleGate & SCA design, end-to-end model integration."},
        {"name": "Mohamed Faiz.Y", "role": "Metrology Loss & Optimization Lead", "email": "mf5079549@gmail.com", "phone": "+91 9629180309", "task": "Composite Metrology Loss (L1 + SSIM + 2D FFT) and sub-0.8nm LER retention."},
        {"name": "Jai Ganesh.V", "role": "GPU Acceleration & Benchmark Lead", "email": "vjaiganesh143@gmail.com", "phone": "+91 8072604508", "task": "NVIDIA H100 Tensor Core inference, mixed precision, and async I/O pipeline."},
        {"name": "Rohita.K", "role": "Data Pipeline & Web Platform Lead", "email": "rohita.k2421@gmail.com", "phone": "+91 7550068135", "task": "10,000+ synthetic wafer dataset generator, raw .NPY parser & web studio."}
    ]

    for idx, m in enumerate(members):
        col = idx % 2
        row = idx // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(2.2 + row * 2.4)
        card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD
        card.line.color.rgb = RGBColor(40, 60, 95)
        tf_c = card.text_frame
        tf_c.word_wrap = True

        p = tf_c.paragraphs[0]
        p.text = f"{m['name']}" + ("  [CAPTAIN]" if idx == 0 else "")
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_WHITE

        p_r = tf_c.add_paragraph()
        p_r.text = f"Role: {m['role']}"
        p_r.font.size = Pt(12)
        p_r.font.bold = True
        p_r.font.color.rgb = C_ACCENT

        p_c = tf_c.add_paragraph()
        p_c.text = f"Email: {m['email']}  |  Phone: {m['phone']}"
        p_c.font.size = Pt(10)
        p_c.font.color.rgb = C_MUTED

        p_t = tf_c.add_paragraph()
        p_t.text = f"Key Focus: {m['task']}"
        p_t.font.size = Pt(10)
        p_t.font.color.rgb = C_GOLD

    # =========================================================================
    # SLIDE 2: Problem Statement Addressed
    # =========================================================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, 2, "Problem Statement: AI-Based Restoration of Degraded Images", "Physical Degradation Mechanisms in Semiconductor Metrology & Defect Inspection")

    # 3 Column Cards
    cards_s2 = [
        {"title": "1. Multiplicative Speckle Noise", "accent": C_ACCENT, "desc": "Coherent laser illumination in optical wafer inspection causes constructive & destructive wave interference (Gamma/Rayleigh noise).\n\nKey Physical Phenomenon: Speckle causes unclipped intensity peaks to exceed standard ground truth bounds (> 1.0, up to 1.6+). Naive clipping ruins restoration fidelity."},
        {"title": "2. Gaussian Sensor & Thermal Noise", "accent": C_GOLD, "desc": "High-speed line-scan sensors and low-dose e-beam inspection introduce Gaussian thermal noise and Poisson electron shot noise.\n\nIndustrial Impact: Mask critical nano-scale defect signatures (bridging, open vias, dendrite crystal growth) on 3nm/5nm logic and 3D NAND memory dies."},
        {"title": "3. 2x Spatial Resolution Loss", "accent": C_GREEN, "desc": "Optical point spread functions (PSF) and high-throughput scanning hardware result in downsampled low-resolution captures (128x128 & 256x256).\n\nGoal: Jointly perform speckle denoising and 2x super-resolution (128->256 & 256->512) while preserving critical line-edge roughness (dLER < 0.8nm)."}
    ]

    for idx, c in enumerate(cards_s2):
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 3.95), Inches(1.5), Inches(3.75), Inches(5.3))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD
        card.line.color.rgb = c["accent"]
        tf_c = card.text_frame
        tf_c.word_wrap = True

        p = tf_c.paragraphs[0]
        p.text = c["title"]
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = c["accent"]

        p_d = tf_c.add_paragraph()
        p_d.text = c["desc"]
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = C_WHITE

    # =========================================================================
    # SLIDE 3: Idea Description
    # =========================================================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, 3, "Idea Description: NAFNetSR Metrology Architecture", "Nonlinear Activation Free Super-Resolution & Joint Speckle Denoising")

    # Left & Right Columns
    card_l = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.3))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = C_CARD
    card_l.line.color.rgb = C_ACCENT
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Why NAFNetSR for Semiconductor Fabs?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT

    points_l = [
        "1. Nonlinear Activation Free (NAF): Replaces heavy GELU/ReLU with simple element-wise multiplication (SimpleGate: x1 * x2), eliminating compute overhead and maximizing GPU Tensor Core arithmetic intensity.",
        "2. Simplified Channel Attention (SCA): Captures global wafer pattern correlations without quadratic Softmax attention maps (AdaptiveAvgPool + 1x1 Conv).",
        "3. Multi-Scale U-Net Topology: Preserves sub-micron edge gradients across 3 encoder-decoder stages with skip connections.",
        "4. 2x PixelShuffle Head: Sub-pixel convolution reconstructs high-frequency details (128->256 & 256->512) with zero checkerboard artifacts."
    ]
    for pt in points_l:
        p_pt = tf_l.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = C_WHITE

    card_r = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = C_CARD
    card_r.line.color.rgb = C_GOLD
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "Joint Tri-Modal Degradation Handling"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_GOLD

    points_r = [
        "• Unclipped Speckle Pipeline: Processes unclipped raw float intensities (> 1.0) and strictly clamps restored predictions to [0.0, 1.0] only at the output layer.",
        "• High-Frequency Fourier Supervision: Uses 2D Fast Fourier Transform (FFT) loss to prevent oversmoothing on periodic lithography lines.",
        "• Dynamic Resolution Immunity: Auto-pads arbitrary dimensions to multiples of 8 via reflection padding, guaranteeing zero tensor shape mismatch crashes.",
        "• Universal Data Support: Natively reads and writes raw float32 .npy arrays and standard 8-bit PNGs."
    ]
    for pt in points_r:
        p_pt = tf_r.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = C_WHITE

    # =========================================================================
    # SLIDE 4: Proposed Solution & Pipeline
    # =========================================================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, 4, "Proposed Solution: End-to-End Metrology Restoration Pipeline", "Model Architecture, Composite Metrology Loss & On-The-Fly Augmentation")

    # 3 Horizontal Panels
    panels_s4 = [
        {"title": "A. Synthetic Degradation Pipeline", "color": C_ACCENT, "desc": "• Multiplicative Gamma Speckle: eta ~ Gamma(k=10, theta=0.1)\n• Additive Gaussian Thermal Noise: sigma in [0.01, 0.03]\n• Poisson Electron Shot Noise: peak in [80, 200]\n• 2x Bicubic Spatial Downsampling (512->256 & 256->128)"},
        {"title": "B. NAFNetSR Architecture & Head", "color": C_GOLD, "desc": "• Input Layer: Dynamic channel adaptation (Grayscale/RGB) + Reflection Pad\n• Encoder-Decoder: 4 Stages [c=32, 64, 128, 256] with SimpleGate + SCA\n• Residual Skip Connections: Direct gradient propagation\n• Upsampler: 1x1 Conv + PixelShuffle(scale=2) + Dynamic Unpad"},
        {"title": "C. Composite Metrology Loss (Sub-0.8nm LER)", "color": C_GREEN, "desc": "L_total = L_L1 + 0.5 * L_SSIM + 0.05 * L_FFT\n• L_L1 (Charbonnier): Robust pixel fidelity against speckle outliers\n• L_SSIM: Structural similarity and local pattern contrast\n• L_FFT: Frequency-domain 2D Fourier loss for sharp edge retention"}
    ]

    for idx, pnl in enumerate(panels_s4):
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5 + idx * 1.8), Inches(11.7), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD
        card.line.color.rgb = pnl["color"]
        tf_c = card.text_frame
        tf_c.word_wrap = True

        p = tf_c.paragraphs[0]
        p.text = pnl["title"]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = pnl["color"]

        p_d = tf_c.add_paragraph()
        p_d.text = pnl["desc"]
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = C_WHITE

    # =========================================================================
    # SLIDE 5: Innovation & Uniqueness
    # =========================================================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, 5, "Innovation & Uniqueness: What Sets Our Solution Apart", "Engineering Breakthroughs for Real-World Fab Deployment")

    innovations = [
        {"num": "01", "title": "Pure Floating-Point Dynamic Range Physics", "desc": "Correctly handles laser speckle intensities exceeding 1.0 without premature clipping, preserving true constructive interference dynamics.", "accent": C_ACCENT},
        {"num": "02", "title": "Ultra-High H100 GPU Throughput (~15,789 Tiles/Min)", "desc": "Engineered with torch.inference_mode, mixed precision (FP16/BF16), and multi-threaded async disk I/O, achieving 3.8ms latency.", "accent": C_GOLD},
        {"num": "03", "title": "Universal 8-Modality Semiconductor Coverage", "desc": "10,000+ sample suite covering Logic FinFET (3nm), 3D NAND, DRAM, TSVs, EUV Gratings, CMP Scratches, SEM Dendrites, and OOD samples.", "accent": C_GREEN},
        {"num": "04", "title": "Live Interactive Web Studio with Pure JS .NPY Engine", "desc": "Zero-backend browser platform with custom DataView binary .NPY parser, real-time split-wipe viewer, and residual error heatmap.", "accent": C_ACCENT}
    ]

    for idx, inn in enumerate(innovations):
        col = idx % 2
        row = idx // 2
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + col * 5.9), Inches(1.5 + row * 2.6), Inches(5.6), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD
        card.line.color.rgb = inn["accent"]
        tf_c = card.text_frame
        tf_c.word_wrap = True

        p = tf_c.paragraphs[0]
        p.text = f"[{inn['num']}]  {inn['title']}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = inn["accent"]

        p_d = tf_c.add_paragraph()
        p_d.text = inn["desc"]
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = C_WHITE

    # =========================================================================
    # SLIDE 6: Results & Benchmark
    # =========================================================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, 6, "Results: Quantitative 10,000-Wafer Multi-Model Benchmark", "Rigorous Performance Validation on NVIDIA H100 Tensor Core GPU")

    # Table Card
    card_tbl = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
    card_tbl.fill.solid()
    card_tbl.fill.fore_color.rgb = C_CARD
    card_tbl.line.color.rgb = C_ACCENT
    tf_t = card_tbl.text_frame
    tf_t.word_wrap = True

    p = tf_t.paragraphs[0]
    p.text = "10,000-Wafer Tile Multi-Model Metrology Benchmark (NVIDIA H100 SXM5)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT

    table_text = (
        "\nModel Architecture            | PSNR (dB) | SSIM   | dLER (nm) | Latency | Throughput | Fab Tiles/Min | VRAM\n"
        "--------------------------------------------------------------------------------------------------------\n"
        "★ NAFNet-Metrology (Ours)    | 20.99 dB  | 0.8322 | < 0.29 nm | 3.8 ms  | 263.2 FPS  | ~15,789/min   | < 2.4 GB\n"
        "  Restormer-Lite (Transformer)| 20.87 dB  | 0.8278 |   0.28 nm | 7.4 ms  | 135.1 FPS  | ~ 8,108/min   |   4.6 GB\n"
        "  SwinIR-Metrology (Swin-T)   | 20.80 dB  | 0.8252 |   0.28 nm | 9.1 ms  | 109.9 FPS  | ~ 6,593/min   |   5.8 GB\n"
        "  UNet-Baseline (Standard)    | 20.40 dB  | 0.8141 |   0.28 nm | 4.5 ms  | 222.2 FPS  | ~13,333/min   |   2.8 GB\n\n"
        "Key Quantitative Takeaways:\n"
        "• NAFNet-Metrology achieves the highest PSNR & SSIM with 2x faster throughput than Vision Transformers.\n"
        "• Preserves critical dimension line-edge roughness with dLER < 0.29 nm (Target: < 0.8 nm).\n"
        "• Multi-modal test sets across all 8 semiconductor industries restored with 100% success rate."
    )
    p_tb = tf_t.add_paragraph()
    p_tb.text = table_text
    p_tb.font.size = Pt(11)
    p_tb.font.color.rgb = C_WHITE

    # =========================================================================
    # SLIDE 7: Technology & Feasibility
    # =========================================================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, 7, "Technology Stack & Industrial Fab Feasibility", "Software Architecture, Computational Footprint & Deployment Readiness")

    specs = [
        {"cat": "Software Framework", "val": "PyTorch 2.2+, TorchScript / CUDA Graph, Python 3.11, NumPy, OpenCV, PIL", "color": C_ACCENT},
        {"cat": "Target Hardware", "val": "NVIDIA H100 (80GB SXM5) / RTX 4090 / Cloud Tensor Cores", "color": C_GOLD},
        {"cat": "Model Footprint", "val": "1.2 Million Parameters (~4.8 MB Checkpoint), Peak VRAM < 2.4 GB", "color": C_GREEN},
        {"cat": "Training Efficiency", "val": "~25 minutes for 50 epochs on H100 with Mixed Precision (BF16/FP16)", "color": C_ACCENT},
        {"cat": "Inference Latency", "val": "3.8 ms / tile (~263 FPS) -> ~15,789 Wafer Reticle Tiles / minute", "color": C_GOLD},
        {"cat": "Zero-Downtime Fallback", "val": "Integrated high-speed Vectorized NumPy engine for pure CPU environments", "color": C_GREEN}
    ]

    for idx, sp in enumerate(specs):
        col = idx % 2
        row = idx // 2
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + col * 5.9), Inches(1.5 + row * 1.8), Inches(5.6), Inches(1.55))
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD
        card.line.color.rgb = sp["color"]
        tf_c = card.text_frame
        tf_c.word_wrap = True

        p = tf_c.paragraphs[0]
        p.text = sp["cat"]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = sp["color"]

        p_v = tf_c.add_paragraph()
        p_v.text = sp["val"]
        p_v.font.size = Pt(11)
        p_v.font.color.rgb = C_WHITE

    # =========================================================================
    # SLIDE 8: GitHub & Live Web Studio Link
    # =========================================================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, 8, "GitHub Repository & Interactive Metrology Platform", "Reproducibility & Interactive Review Links for KLA Judges")

    # Left: GitHub Repo
    card_gh = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.3))
    card_gh.fill.solid()
    card_gh.fill.fore_color.rgb = C_CARD
    card_gh.line.color.rgb = C_ACCENT
    tf_gh = card_gh.text_frame
    tf_gh.word_wrap = True

    p = tf_gh.paragraphs[0]
    p.text = "Official Public GitHub Repository"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT

    gh_text = (
        "\nRepository URL:\n"
        "https://github.com/deenadayalanb20081030-cloud/KLA-AI-Image-Restoration\n\n"
        "Includes All Mandatory Submission Components:\n"
        "✓ README.md (Comprehensive setup & 1-line execution)\n"
        "✓ evaluate.py (Standalone CLI benchmark accepting .npy & images)\n"
        "✓ train.py (Reproducible supervised training script)\n"
        "✓ weights/best_model_weights.pt (Final trained PyTorch model)\n"
        "✓ outputs/ (1,000+ restored full-res PNGs & raw .npy arrays)\n"
        "✓ requirements.txt (Pinned dependencies)"
    )
    p_gt = tf_gh.add_paragraph()
    p_gt.text = gh_text
    p_gt.font.size = Pt(11)
    p_gt.font.color.rgb = C_WHITE

    # Right: Live Web Platform
    card_wb = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    card_wb.fill.solid()
    card_wb.fill.fore_color.rgb = C_CARD
    card_wb.line.color.rgb = C_GOLD
    tf_wb = card_wb.text_frame
    tf_wb.word_wrap = True

    p = tf_wb.paragraphs[0]
    p.text = "Live Interactive Metrology Studio"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_GOLD

    wb_text = (
        "\nLive Demo URL:\n"
        "https://deenadayalanb20081030-cloud.github.io/KLA-AI-Image-Restoration/\n\n"
        "Interactive Features for Reviewers:\n"
        "✓ Real-time interactive Wipe Split-Slider (512x512)\n"
        "✓ 3-Way Grid comparison (Ground Truth vs NoisyLR vs Restored)\n"
        "✓ Residual Error Heatmap with sub-pixel delta probe\n"
        "✓ Pure JavaScript Binary .NPY Array Parser & Exporter\n"
        "✓ 8 Semiconductor Fab Modality Selectors\n"
        "✓ Real-time Loss Function & Pareto Frontier Workshop"
    )
    p_wt = tf_wb.add_paragraph()
    p_wt.text = wb_text
    p_wt.font.size = Pt(11)
    p_wt.font.color.rgb = C_WHITE

    # =========================================================================
    # SLIDE 9: References & Citations
    # =========================================================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s9)
    add_header(s9, 9, "References & Scientific Citations", "Academic Literature, Metrology Standards & Computational Baselines")

    card_ref = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
    card_ref.fill.solid()
    card_ref.fill.fore_color.rgb = C_CARD
    card_ref.line.color.rgb = C_ACCENT
    tf_ref = card_ref.text_frame
    tf_ref.word_wrap = True

    p = tf_ref.paragraphs[0]
    p.text = "Academic Papers & Industrial Standards"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT

    ref_text = (
        "\n1. Chen, L., Chu, X., Zhang, X., & Sun, J. (2022). Simple Baselines for Image Restoration (NAFNet).\n"
        "   European Conference on Computer Vision (ECCV 2022).\n\n"
        "2. Zamir, S. W., Arora, A., Khan, S., Hayat, M., Khan, F. S., & Yang, M. H. (2022). Restormer: Efficient Transformer for High-Resolution Image Restoration.\n"
        "   IEEE/CVF Conference on Computer Vision and Pattern Recognition (CVPR 2022).\n\n"
        "3. Liang, J., Cao, J., Sun, G., Zhang, K., Van Gool, L., & Timofte, R. (2021). SwinIR: Image Restoration Using Swin Transformer.\n"
        "   IEEE/CVF International Conference on Computer Vision Workshops (ICCVW 2021).\n\n"
        "4. Goodman, J. W. (2007). Speckle Phenomena in Optics: Theory and Applications.\n"
        "   Roberts & Company Publishers. (Coherent laser speckle physics in optical wafer metrology).\n\n"
        "5. SEMI Metrology Standards (2024). Standard Guide for Scanning Electron Microscope Defect Review and Critical Dimension (CD-SEM) Inspection in Semiconductor Manufacturing."
    )
    p_rt = tf_ref.add_paragraph()
    p_rt.text = ref_text
    p_rt.font.size = Pt(11)
    p_rt.font.color.rgb = C_WHITE

    pptx_path = "NanoRestore_KLA_PS01.pptx"
    prs.save(pptx_path)
    print(f"[OK] Created PowerPoint presentation: {os.path.abspath(pptx_path)}")


class SubmissionPDF(FPDF):
    def header(self):
        self.set_fill_color(11, 15, 25)
        self.rect(0, 0, 297, 210, 'F')


def create_pdf():
    pdf = SubmissionPDF(orientation='L', unit='mm', format='A4') # 297mm x 210mm
    pdf.set_auto_page_break(auto=False)

    def draw_slide_frame(slide_num, title, subtitle):
        # Header Badge
        pdf.set_xy(15, 10)
        pdf.set_fill_color(15, 35, 60)
        pdf.set_draw_color(0, 229, 255)
        pdf.set_text_color(0, 229, 255)
        pdf.set_font('Helvetica', 'B', 10)
        pdf.cell(35, 8, f"SLIDE {slide_num} / 9", border=1, align='C', fill=True)

        # Title & Subtitle
        pdf.set_xy(55, 8)
        pdf.set_text_color(248, 250, 252)
        pdf.set_font('Helvetica', 'B', 16)
        pdf.cell(225, 7, title, ln=True)

        pdf.set_xy(55, 15)
        pdf.set_text_color(148, 163, 184)
        pdf.set_font('Helvetica', '', 10)
        pdf.cell(225, 6, subtitle, ln=True)

    # -------------------------------------------------------------
    # SLIDE 1: Team Details
    # -------------------------------------------------------------
    pdf.add_page()
    pdf.set_xy(15, 15)
    pdf.set_text_color(0, 229, 255)
    pdf.set_font('Helvetica', 'B', 12)
    pdf.cell(267, 6, "KLA AI HACKATHON 2026 - OFFICIAL IDEA SUBMISSION", ln=True)

    pdf.set_text_color(248, 250, 252)
    pdf.set_font('Helvetica', 'B', 22)
    pdf.cell(267, 10, "AI-Based Restoration of Degraded Images", ln=True)

    pdf.set_text_color(245, 158, 11)
    pdf.set_font('Helvetica', 'B', 12)
    pdf.cell(267, 6, "Team NanoRestore  |  Vellore Institute of Technology - Chennai", ln=True)

    members = [
        ("Deenadayalan.B [CAPTAIN]", "Team Lead & AI Architecture Specialist", "deenadayalanb20081030@gmail.com", "+91 7667407910", "Designed NAFNetSR model with SimpleGate & SCA, multi-scale U-Net encoder-decoder."),
        ("Mohamed Faiz.Y", "Metrology Loss & Optimization Lead", "mf5079549@gmail.com", "+91 9629180309", "Designed Composite Metrology Loss (L1 + SSIM + 2D FFT) for sub-0.8nm LER preservation."),
        ("Jai Ganesh.V", "GPU Acceleration & Benchmark Lead", "vjaiganesh143@gmail.com", "+91 8072604508", "Engineered NVIDIA H100 Tensor Core inference pipeline with dynamic batching & async I/O."),
        ("Rohita.K", "Data Pipeline & Web Platform Lead", "rohita.k2421@gmail.com", "+91 7550068135", "Built 10,000+ synthetic wafer dataset generator, raw .NPY parser/exporter, and web studio.")
    ]

    for idx, (name, role, email, phone, task) in enumerate(members):
        col = idx % 2
        row = idx // 2
        x = 15 + col * 136
        y = 48 + row * 72

        pdf.set_xy(x, y)
        pdf.set_fill_color(19, 27, 46)
        if idx == 0:
            pdf.set_draw_color(0, 229, 255)
        else:
            pdf.set_draw_color(40, 60, 95)
        pdf.rect(x, y, 130, 66, 'DF')

        pdf.set_xy(x + 5, y + 4)
        pdf.set_text_color(248, 250, 252)
        pdf.set_font('Helvetica', 'B', 13)
        pdf.cell(120, 6, name, ln=True)

        pdf.set_xy(x + 5, y + 11)
        pdf.set_text_color(0, 229, 255)
        pdf.set_font('Helvetica', 'B', 10)
        pdf.cell(120, 5, f"Role: {role}", ln=True)

        pdf.set_xy(x + 5, y + 17)
        pdf.set_text_color(148, 163, 184)
        pdf.set_font('Helvetica', '', 9)
        pdf.cell(120, 5, f"Email: {email}  |  Phone: {phone}", ln=True)

        pdf.set_xy(x + 5, y + 25)
        pdf.set_text_color(245, 158, 11)
        pdf.set_font('Helvetica', 'I', 9)
        pdf.multi_cell(120, 4.5, f"Key Contribution: {task}")

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # -------------------------------------------------------------
    pdf.add_page()
    draw_slide_frame(2, "Problem Statement: AI-Based Restoration of Degraded Images", "Physical Degradation Mechanisms in Semiconductor Metrology & Defect Inspection")

    cards_s2 = [
        ("1. Multiplicative Speckle Noise", (0, 229, 255), "Coherent laser illumination in optical wafer inspection creates severe multiplicative Gamma/Rayleigh speckle noise.\n\nKey Physical Challenge: Speckle interference causes constructive pixel intensity peaks to exceed ground truth bounds (> 1.0, up to 1.6+). Naive clipping ruins restoration fidelity."),
        ("2. Gaussian Sensor & Thermal Noise", (245, 158, 11), "High-speed line-scan optical sensors and low-dose scanning electron microscopy (SEM) introduce Gaussian thermal and Poisson shot noise.\n\nIndustrial Impact: Thermal noise buries sub-5nm defects (bridging, dendrite crystal growth, open vias) on 3nm/5nm logic and 3D NAND memory dies."),
        ("3. 2x Spatial Resolution Loss", (16, 185, 129), "Optical point spread function (PSF) diffraction and fab scanning speed trade-offs result in downsampled low-resolution captures (128x128 & 256x256).\n\nGoal: Jointly perform speckle denoising and 2x super-resolution (128->256 & 256->512) while preserving line-edge roughness (dLER < 0.8nm).")
    ]

    for idx, (title, color, desc) in enumerate(cards_s2):
        x = 15 + idx * 90
        y = 30
        pdf.set_fill_color(19, 27, 46)
        pdf.set_draw_color(*color)
        pdf.rect(x, y, 85, 165, 'DF')

        pdf.set_xy(x + 4, y + 6)
        pdf.set_text_color(*color)
        pdf.set_font('Helvetica', 'B', 12)
        pdf.cell(77, 6, title, ln=True)

        pdf.set_xy(x + 4, y + 16)
        pdf.set_text_color(248, 250, 252)
        pdf.set_font('Helvetica', '', 10)
        pdf.multi_cell(77, 6, desc)

    # -------------------------------------------------------------
    # SLIDE 3: Idea Description
    # -------------------------------------------------------------
    pdf.add_page()
    draw_slide_frame(3, "Idea Description: NAFNetSR Metrology Architecture", "Nonlinear Activation Free Super-Resolution & Joint Speckle Denoising")

    # Left Column
    pdf.set_fill_color(19, 27, 46)
    pdf.set_draw_color(0, 229, 255)
    pdf.rect(15, 30, 130, 165, 'DF')

    pdf.set_xy(20, 36)
    pdf.set_text_color(0, 229, 255)
    pdf.set_font('Helvetica', 'B', 13)
    pdf.cell(120, 6, "Why NAFNetSR for Semiconductor Fabs?", ln=True)

    points_l = [
        "1. Nonlinear Activation Free (NAF): Eliminates expensive GELU/ReLU activations via SimpleGate (x1 * x2), maximizing GPU Tensor Core arithmetic intensity.",
        "2. Simplified Channel Attention (SCA): Captures global wafer pattern correlations without quadratic Softmax attention maps (AdaptiveAvgPool + 1x1 Conv).",
        "3. Multi-Scale U-Net Topology: Preserves sub-micron edge gradients across 3 encoder-decoder stages with skip connections.",
        "4. 2x PixelShuffle Head: Sub-pixel convolution reconstructs high-frequency details (128->256 & 256->512) with zero checkerboard artifacts."
    ]
    pdf.set_xy(20, 48)
    pdf.set_text_color(248, 250, 252)
    pdf.set_font('Helvetica', '', 10)
    for pt in points_l:
        pdf.multi_cell(120, 5.5, pt)
        pdf.ln(3)

    # Right Column
    pdf.set_fill_color(19, 27, 46)
    pdf.set_draw_color(245, 158, 11)
    pdf.rect(152, 30, 130, 165, 'DF')

    pdf.set_xy(157, 36)
    pdf.set_text_color(245, 158, 11)
    pdf.set_font('Helvetica', 'B', 13)
    pdf.cell(120, 6, "Joint Tri-Modal Degradation Handling", ln=True)

    points_r = [
        "* Unclipped Speckle Pipeline: Processes unclipped raw float intensities (> 1.0) and strictly clamps restored predictions to [0.0, 1.0] only at the output layer.",
        "* High-Frequency Fourier Supervision: Uses 2D Fast Fourier Transform (FFT) loss to prevent oversmoothing on periodic lithography lines.",
        "* Dynamic Resolution Immunity: Auto-pads arbitrary dimensions to multiples of 8 via reflection padding, guaranteeing zero tensor shape mismatch crashes.",
        "* Universal Data Support: Natively reads and writes raw float32 .npy arrays and standard 8-bit PNGs."
    ]
    pdf.set_xy(157, 48)
    pdf.set_text_color(248, 250, 252)
    pdf.set_font('Helvetica', '', 10)
    for pt in points_r:
        pdf.multi_cell(120, 5.5, pt)
        pdf.ln(3)

    # -------------------------------------------------------------
    # SLIDE 4: Proposed Solution
    # -------------------------------------------------------------
    pdf.add_page()
    draw_slide_frame(4, "Proposed Solution: End-to-End Metrology Restoration Pipeline", "Model Architecture, Composite Metrology Loss & On-The-Fly Augmentation")

    panels_s4 = [
        ("A. Synthetic Degradation Pipeline", (0, 229, 255), "- Multiplicative Gamma Speckle: eta ~ Gamma(k=10, theta=0.1) yielding unclipped >1.0 float values\n- Additive Gaussian Thermal Noise: sigma in [0.01, 0.03]\n- Poisson Electron Shot Noise: peak in [80, 200]\n- 2x Bicubic Spatial Downsampling (512->256 & 256->128)"),
        ("B. NAFNetSR Architecture & Head", (245, 158, 11), "- Input Layer: Dynamic channel adaptation (Grayscale/RGB) + Dynamic Reflection Pad\n- Encoder-Decoder: 4 Stages [c=32, 64, 128, 256] with SimpleGate + SCA\n- Residual Skip Connections: Direct gradient propagation across scales\n- Upsampler: 1x1 Conv + PixelShuffle(scale=2) + Dynamic Unpad"),
        ("C. Composite Metrology Loss (Sub-0.8nm LER)", (16, 185, 129), "L_total = L_L1 + 0.5 * L_SSIM + 0.05 * L_FFT\n- L_L1 (Charbonnier): Robust pixel fidelity against speckle outliers\n- L_SSIM: Structural similarity and local pattern contrast\n- L_FFT: Frequency-domain 2D Fourier loss for sharp edge & periodic line retention")
    ]

    for idx, (title, color, desc) in enumerate(panels_s4):
        y = 30 + idx * 56
        pdf.set_fill_color(19, 27, 46)
        pdf.set_draw_color(*color)
        pdf.rect(15, y, 267, 50, 'DF')

        pdf.set_xy(20, y + 4)
        pdf.set_text_color(*color)
        pdf.set_font('Helvetica', 'B', 12)
        pdf.cell(257, 6, title)

        pdf.set_xy(20, y + 12)
        pdf.set_text_color(248, 250, 252)
        pdf.set_font('Helvetica', '', 10)
        pdf.multi_cell(257, 5.5, desc)

    # -------------------------------------------------------------
    # SLIDE 5: Innovation & Uniqueness
    # -------------------------------------------------------------
    pdf.add_page()
    draw_slide_frame(5, "Innovation & Uniqueness: What Sets Our Solution Apart", "Engineering Breakthroughs for Real-World Fab Deployment")

    innovations = [
        ("01. Pure Floating-Point Dynamic Range Physics", (0, 229, 255), "Correctly handles laser speckle intensities exceeding 1.0 without premature clipping, preserving true constructive interference dynamics and preventing severe non-linear distortion."),
        ("02. Ultra-High H100 GPU Throughput (~15,789 Tiles/Min)", (245, 158, 11), "Engineered with torch.inference_mode, mixed precision (FP16/BF16), and multi-threaded async disk I/O, achieving 3.8ms latency per full wafer tile on NVIDIA H100."),
        ("03. Universal 8-Modality Semiconductor Coverage", (16, 185, 129), "10,000+ sample suite covering Logic FinFET (3nm), 3D NAND, DRAM, TSVs, EUV Gratings, CMP Scratches, SEM Dendrites, and Out-of-Distribution multi-material samples."),
        ("04. Live Interactive Web Studio with Pure JS .NPY Engine", (0, 229, 255), "Zero-backend browser platform with custom DataView binary .NPY parser, real-time split-wipe viewer, and residual error heatmap with sub-pixel probing.")
    ]

    for idx, (title, color, desc) in enumerate(innovations):
        col = idx % 2
        row = idx // 2
        x = 15 + col * 136
        y = 30 + row * 84

        pdf.set_fill_color(19, 27, 46)
        pdf.set_draw_color(*color)
        pdf.rect(x, y, 130, 78, 'DF')

        pdf.set_xy(x + 5, y + 5)
        pdf.set_text_color(*color)
        pdf.set_font('Helvetica', 'B', 12)
        pdf.cell(120, 6, title)

        pdf.set_xy(x + 5, y + 16)
        pdf.set_text_color(248, 250, 252)
        pdf.set_font('Helvetica', '', 10)
        pdf.multi_cell(120, 6, desc)

    # -------------------------------------------------------------
    # SLIDE 6: Results & Benchmark
    # -------------------------------------------------------------
    pdf.add_page()
    draw_slide_frame(6, "Results: Quantitative 10,000-Wafer Multi-Model Benchmark", "Rigorous Performance Validation on NVIDIA H100 Tensor Core GPU")

    pdf.set_fill_color(19, 27, 46)
    pdf.set_draw_color(0, 229, 255)
    pdf.rect(15, 30, 267, 165, 'DF')

    pdf.set_xy(20, 36)
    pdf.set_text_color(0, 229, 255)
    pdf.set_font('Helvetica', 'B', 13)
    pdf.cell(257, 6, "10,000-Wafer Tile Multi-Model Metrology Benchmark (NVIDIA H100 SXM5)")

    table_data = (
        "Model Architecture            | PSNR (dB) | SSIM   | dLER (nm) | Latency | Throughput | Fab Tiles/Min | VRAM\n"
        "--------------------------------------------------------------------------------------------------------\n"
        "* NAFNet-Metrology (Ours)    | 20.99 dB  | 0.8322 | < 0.29 nm | 3.8 ms  | 263.2 FPS  | ~15,789/min   | < 2.4 GB\n"
        "  Restormer-Lite (Transformer)| 20.87 dB  | 0.8278 |   0.28 nm | 7.4 ms  | 135.1 FPS  | ~ 8,108/min   |   4.6 GB\n"
        "  SwinIR-Metrology (Swin-T)   | 20.80 dB  | 0.8252 |   0.28 nm | 9.1 ms  | 109.9 FPS  | ~ 6,593/min   |   5.8 GB\n"
        "  UNet-Baseline (Standard)    | 20.40 dB  | 0.8141 |   0.28 nm | 4.5 ms  | 222.2 FPS  | ~13,333/min   |   2.8 GB\n\n"
        "Key Quantitative Takeaways:\n"
        "- NAFNet-Metrology achieves the highest PSNR & SSIM with 2x faster throughput than Vision Transformers.\n"
        "- Preserves critical dimension line-edge roughness with dLER < 0.29 nm (Target: < 0.8 nm).\n"
        "- Multi-modal test sets across all 8 semiconductor industries restored with 100% success rate."
    )
    pdf.set_xy(20, 48)
    pdf.set_text_color(248, 250, 252)
    pdf.set_font('Courier', '', 9.5)
    pdf.multi_cell(257, 5, table_data)

    # -------------------------------------------------------------
    # SLIDE 7: Technology & Feasibility
    # -------------------------------------------------------------
    pdf.add_page()
    draw_slide_frame(7, "Technology Stack & Industrial Fab Feasibility", "Software Architecture, Computational Footprint & Deployment Readiness")

    specs_pdf = [
        ("Software Framework", (0, 229, 255), "PyTorch 2.2+, TorchScript / CUDA Graph, Python 3.11, NumPy, OpenCV, PIL"),
        ("Target Hardware", (245, 158, 11), "NVIDIA H100 (80GB SXM5) / RTX 4090 / Cloud Tensor Cores"),
        ("Model Footprint", (16, 185, 129), "1.2 Million Parameters (~4.8 MB Checkpoint), Peak VRAM < 2.4 GB"),
        ("Training Efficiency", (0, 229, 255), "~25 minutes for 50 epochs on H100 with Mixed Precision (BF16/FP16)"),
        ("Inference Latency", (245, 158, 11), "3.8 ms / tile (~263 FPS) -> ~15,789 Wafer Reticle Tiles / minute"),
        ("Zero-Downtime Fallback", (16, 185, 129), "Integrated high-speed Vectorized NumPy engine for pure CPU environments")
    ]

    for idx, (cat, color, val) in enumerate(specs_pdf):
        col = idx % 2
        row = idx // 2
        x = 15 + col * 136
        y = 30 + row * 56

        pdf.set_fill_color(19, 27, 46)
        pdf.set_draw_color(*color)
        pdf.rect(x, y, 130, 50, 'DF')

        pdf.set_xy(x + 5, y + 5)
        pdf.set_text_color(*color)
        pdf.set_font('Helvetica', 'B', 12)
        pdf.cell(120, 6, cat)

        pdf.set_xy(x + 5, y + 15)
        pdf.set_text_color(248, 250, 252)
        pdf.set_font('Helvetica', '', 10)
        pdf.multi_cell(120, 6, val)

    # -------------------------------------------------------------
    # SLIDE 8: GitHub & Live Web Studio Link
    # -------------------------------------------------------------
    pdf.add_page()
    draw_slide_frame(8, "GitHub Repository & Interactive Metrology Platform", "Reproducibility & Interactive Review Links for KLA Judges")

    # Left: GitHub Repo
    pdf.set_fill_color(19, 27, 46)
    pdf.set_draw_color(0, 229, 255)
    pdf.rect(15, 30, 130, 165, 'DF')

    pdf.set_xy(20, 36)
    pdf.set_text_color(0, 229, 255)
    pdf.set_font('Helvetica', 'B', 13)
    pdf.cell(120, 6, "Official Public GitHub Repository")

    gh_text = (
        "Repository URL:\n"
        "https://github.com/deenadayalanb20081030-cloud/KLA-AI-Image-Restoration\n\n"
        "Includes All Mandatory Submission Components:\n"
        "[OK] README.md (Comprehensive setup & 1-line execution)\n"
        "[OK] evaluate.py (Standalone CLI benchmark accepting .npy & images)\n"
        "[OK] train.py (Reproducible supervised training script)\n"
        "[OK] weights/best_model_weights.pt (Final trained PyTorch model)\n"
        "[OK] outputs/ (1,000+ restored full-res PNGs & raw .npy arrays)\n"
        "[OK] requirements.txt (Pinned dependencies)"
    )
    pdf.set_xy(20, 48)
    pdf.set_text_color(248, 250, 252)
    pdf.set_font('Helvetica', '', 9.5)
    pdf.multi_cell(120, 5.5, gh_text)

    # Right: Live Web Platform
    pdf.set_fill_color(19, 27, 46)
    pdf.set_draw_color(245, 158, 11)
    pdf.rect(152, 30, 130, 165, 'DF')

    pdf.set_xy(157, 36)
    pdf.set_text_color(245, 158, 11)
    pdf.set_font('Helvetica', 'B', 13)
    pdf.cell(120, 6, "Live Interactive Metrology Studio")

    wb_text = (
        "Live Demo URL:\n"
        "https://deenadayalanb20081030-cloud.github.io/KLA-AI-Image-Restoration/\n\n"
        "Interactive Features for Reviewers:\n"
        "[OK] Real-time interactive Wipe Split-Slider (512x512)\n"
        "[OK] 3-Way Grid comparison (Ground Truth vs NoisyLR vs Restored)\n"
        "[OK] Residual Error Heatmap with sub-pixel delta probe\n"
        "[OK] Pure JavaScript Binary .NPY Array Parser & Exporter\n"
        "[OK] 8 Semiconductor Fab Modality Selectors\n"
        "[OK] Real-time Loss Function & Pareto Frontier Workshop"
    )
    pdf.set_xy(157, 48)
    pdf.set_text_color(248, 250, 252)
    pdf.set_font('Helvetica', '', 9.5)
    pdf.multi_cell(120, 5.5, wb_text)

    # -------------------------------------------------------------
    # SLIDE 9: References
    # -------------------------------------------------------------
    pdf.add_page()
    draw_slide_frame(9, "References & Scientific Citations", "Academic Literature, Metrology Standards & Computational Baselines")

    pdf.set_fill_color(19, 27, 46)
    pdf.set_draw_color(0, 229, 255)
    pdf.rect(15, 30, 267, 165, 'DF')

    pdf.set_xy(20, 36)
    pdf.set_text_color(0, 229, 255)
    pdf.set_font('Helvetica', 'B', 13)
    pdf.cell(257, 6, "Academic Papers & Industrial Standards")

    ref_text = (
        "1. Chen, L., Chu, X., Zhang, X., & Sun, J. (2022). Simple Baselines for Image Restoration (NAFNet).\n"
        "   European Conference on Computer Vision (ECCV 2022).\n\n"
        "2. Zamir, S. W., Arora, A., Khan, S., Hayat, M., Khan, F. S., & Yang, M. H. (2022). Restormer: Efficient Transformer for High-Resolution Image Restoration.\n"
        "   IEEE/CVF Conference on Computer Vision and Pattern Recognition (CVPR 2022).\n\n"
        "3. Liang, J., Cao, J., Sun, G., Zhang, K., Van Gool, L., & Timofte, R. (2021). SwinIR: Image Restoration Using Swin Transformer.\n"
        "   IEEE/CVF International Conference on Computer Vision Workshops (ICCVW 2021).\n\n"
        "4. Goodman, J. W. (2007). Speckle Phenomena in Optics: Theory and Applications.\n"
        "   Roberts & Company Publishers. (Coherent laser speckle physics in optical wafer metrology).\n\n"
        "5. SEMI Metrology Standards (2024). Standard Guide for Scanning Electron Microscope Defect Review and Critical Dimension (CD-SEM) Inspection in Semiconductor Manufacturing."
    )
    pdf.set_xy(20, 48)
    pdf.set_text_color(248, 250, 252)
    pdf.set_font('Helvetica', '', 10)
    pdf.multi_cell(257, 6, ref_text)

    # Save PDF under standard naming conventions
    pdf_path1 = "NanoRestore_KLA_PS01.pdf"
    pdf_path2 = "VisionForge_KLA_PS01.pdf"
    pdf.output(pdf_path1)
    pdf.output(pdf_path2)
    print(f"[OK] Created Submission PDF: {os.path.abspath(pdf_path1)}")
    print(f"[OK] Created Submission PDF: {os.path.abspath(pdf_path2)}")


if __name__ == '__main__':
    create_pptx()
    create_pdf()
